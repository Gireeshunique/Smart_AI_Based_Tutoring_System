import pythoncom
from docx2pdf import convert
from pptx import Presentation
from reportlab.pdfgen import canvas


def docx_to_pdf(input_path, output_path):
    pythoncom.CoInitialize()
    try:
        convert(input_path, output_path)
    finally:
        pythoncom.CoUninitialize()

def pptx_to_pdf(input_path, output_path):
    prs = Presentation(input_path)
    c = canvas.Canvas(output_path)

    for slide in prs.slides:
        text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

        c.drawString(40, 800, text[:3000])
        c.showPage()

    c.save()
